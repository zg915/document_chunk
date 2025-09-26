"""
Document processing module for converting documents to markdown and chunking.

This module provides functions to:
1. Convert PDF/image files to markdown using the Marker API
2. Chunk markdown content into smaller, manageable pieces for vector storage
3. Save documents and chunks to Weaviate vector database
"""
#!/usr/bin/env python3
# Standard library imports
import logging
import mimetypes
import os
import tempfile
import time
import warnings
import asyncio
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Union


# Third-party imports
# GPU optimizations - set before importing torch-based libraries
os.environ["OMP_NUM_THREADS"] = "1"
try:
    import torch
    if torch.cuda.is_available():
        torch.set_num_threads(1)
        torch.backends.cudnn.benchmark = True
        torch.backends.cuda.matmul.allow_tf32 = True
        torch.backends.cudnn.allow_tf32 = True
        print(f"GPU optimizations enabled: {torch.cuda.get_device_name(0)}")
except Exception as e:
    print(f"GPU optimizations failed, continuing without: {e}")
    torch = None


import requests
import tiktoken
from dotenv import load_dotenv
from PIL import Image
from weaviate import WeaviateClient
from weaviate.auth import AuthApiKey
from weaviate.connect import ConnectionParams, ProtocolParams
from weaviate.util import generate_uuid5

# Unstructured imports
from unstructured.chunking.basic import chunk_elements
from unstructured.chunking.title import chunk_by_title
from unstructured.partition.md import partition_md

# LangChain imports
from langchain_community.document_loaders import (
    UnstructuredPDFLoader,
    UnstructuredWordDocumentLoader,
    UnstructuredExcelLoader,
    UnstructuredPowerPointLoader,
    UnstructuredHTMLLoader,
    UnstructuredCSVLoader,
    UnstructuredEPubLoader,
    UnstructuredODTLoader,
    UnstructuredRTFLoader,
    PyMuPDFLoader
)

# Configuration
warnings.filterwarnings("ignore")
logger = logging.getLogger(__name__)

# Load environment variables
load_dotenv('.env')


class WebhookManager:
    """Manages webhook-based document processing requests and callbacks."""

    def __init__(self):
        """Initialize the webhook manager."""
        self.pending_requests = {}
        self.logger = logging.getLogger(f"{__name__}.WebhookManager")

    def create_request(self, request_id: str, webhook_url: str) -> asyncio.Future:
        """
        Create a new webhook request and return a future for the result.

        Args:
            request_id: Unique identifier for the request
            webhook_url: URL for webhook callbacks

        Returns:
            asyncio.Future object that will be resolved when webhook callback is received
        """
        future = asyncio.Future()

        self.pending_requests[request_id] = {
            "request_id": request_id,
            "webhook_url": webhook_url,
            "future": future,
            "created_at": time.time()
        }

        self.logger.info(f"Created webhook request: {request_id}")
        return future

    def process_callback(self, request_id: str, success: bool,
                         markdown_content: Optional[str] = None,
                         extracted_data: Optional[Dict] = None,
                         error_message: Optional[str] = None) -> bool:
        """
        Process an incoming webhook callback.

        Args:
            request_id: The request ID from the callback
            success: Whether the processing was successful
            markdown_content: The converted markdown content (if successful)
            extracted_data: Additional extracted data (if any)
            error_message: Error message (if failed)

        Returns:
            True if callback was processed, False if request_id not found
        """
        if request_id not in self.pending_requests:
            self.logger.warning(f"Received callback for unknown request_id: {request_id}")
            return False

        webhook_request = self.pending_requests[request_id]
        future = webhook_request["future"]

        try:
            if success and markdown_content:
                result = {
                    "success": True,
                    "markdown_content": markdown_content,
                    "extracted_data": extracted_data or {}
                }
                future.set_result(result)
                self.logger.info(f"Successfully processed webhook callback for request_id: {request_id}")
            else:
                result = {
                    "success": False,
                    "error": error_message or "Webhook callback indicated failure"
                }
                future.set_result(result)
                self.logger.warning(f"Webhook callback failed for request_id: {request_id}")

            # Clean up the pending request
            del self.pending_requests[request_id]
            return True

        except Exception as e:
            self.logger.error(f"Error processing webhook callback for {request_id}: {e}")
            if not future.done():
                future.set_exception(e)
            return False

    def cleanup_timeout_requests(self, timeout_seconds: int = 300):
        """
        Clean up requests that have timed out.

        Args:
            timeout_seconds: Maximum age of requests before cleanup (default: 5 minutes)
        """
        current_time = time.time()
        timeout_requests = []

        for request_id, request_data in self.pending_requests.items():
            if current_time - request_data["created_at"] > timeout_seconds:
                timeout_requests.append(request_id)

        for request_id in timeout_requests:
            request_data = self.pending_requests[request_id]
            future = request_data["future"]

            if not future.done():
                future.set_exception(asyncio.TimeoutError(f"Webhook request {request_id} timed out"))

            del self.pending_requests[request_id]
            self.logger.warning(f"Cleaned up timed out request: {request_id}")

        if timeout_requests:
            self.logger.info(f"Cleaned up {len(timeout_requests)} timed out webhook requests")

    def get_pending_count(self) -> int:
        """Get the number of pending webhook requests."""
        return len(self.pending_requests)


# Global webhook manager instance
_webhook_manager = WebhookManager()


def get_webhook_manager() -> WebhookManager:
    """Get the global webhook manager instance."""
    return _webhook_manager


# Configuration
class Config:
    """Configuration settings for the document processor."""
    API_KEY = os.getenv("MARKER_API_KEY")
    API_URL = os.getenv("MARKER_API_URL")

    
    MARKER_PARAMS = {
        'output_format': 'markdown',
        'use_llm': True,
        'disable_image_extraction': True
    }
    
    SUPPORTED_IMAGE_FORMATS = ['.jpg', '.jpeg', '.png', '.gif', '.tiff', '.webp']
    SUPPORTED_WORD_FORMATS = ['.docx', '.doc']
    SUPPORTED_EXCEL_FORMATS = ['.xlsx', '.xls']
    SUPPORTED_POWERPOINT_FORMATS = ['.pptx', '.ppt']
    SUPPORTED_TEXT_FORMATS = ['.txt', '.md', '.csv', '.html', '.htm']
    SUPPORTED_DOCUMENT_FORMATS = ['.pdf'] + SUPPORTED_IMAGE_FORMATS + SUPPORTED_WORD_FORMATS + SUPPORTED_EXCEL_FORMATS + SUPPORTED_POWERPOINT_FORMATS + SUPPORTED_TEXT_FORMATS
    
    # Chunking parameters
    MAX_CHUNK_SIZE = 2500
    NEW_CHUNK_AFTER = 2000
    MIN_CHUNK_SIZE = 500
    CHUNK_OVERLAP = 200
    
    # API settings
    MAX_RETRIES = 30
    RETRY_DELAY = 10
    
    # Local Marker settings - using marker instead of marker_single for batch processing
    LOCAL_MARKER_COMMAND = "marker"  # Use marker for consistent batch processing
    LOCAL_MARKER_OUTPUT_DIR = os.getenv("LOCAL_MARKER_OUTPUT_DIR", "./marker_output")
    
    # GPU optimization settings
    NUM_WORKERS = int(os.getenv("MARKER_NUM_WORKERS", "4"))  # 4-7 range recommended
    PAGE_BATCH_SIZE = int(os.getenv("MARKER_PAGE_BATCH", "6"))  # 4-8 range recommended
    
    # Weaviate settings
    WEAVIATE_API_KEY = os.getenv("WEAVIATE_API_KEY")
    WEAVIATE_URL = os.getenv("WEAVIATE_URL")
    WEAVIATE_HTTP_PORT = 8080
    WEAVIATE_GRPC_PORT = 50051


class DocumentProcessor:
    """Handles document conversion and processing operations."""
    
    def __init__(self, api_key: Optional[str] = None, api_url: Optional[str] = None):
        """Initialize the document processor."""
        self.api_key = api_key or Config.API_KEY
        self.api_url = api_url or Config.API_URL
        self.headers = {"X-API-Key": self.api_key}
    
    def _convert_image_to_pdf(self, image_path: str) -> str:
        """
        Convert an image file to PDF format.
        
        Args:
            image_path: Path to the image file
            
        Returns:
            Path to the temporary PDF file
        """
        try:
            with Image.open(image_path) as img:
                temp_pdf = tempfile.NamedTemporaryFile(suffix='.pdf', delete=False)
                temp_pdf_path = temp_pdf.name
                temp_pdf.close()
                
                if img.mode != 'RGB':
                    img = img.convert('RGB')
                
                img.save(temp_pdf_path, 'PDF', resolution=150, quality=95)
                logger.info(f"Converted image to PDF: {temp_pdf_path}")
                return temp_pdf_path
        except Exception as e:
            logger.error(f"Failed to convert image to PDF: {e}")
            raise
    
    def _process_text_file(self, file_path: str) -> Optional[str]:
        """
        Process text-based files (txt, md, csv, html).

        Args:
            file_path: Path to the text file

        Returns:
            Markdown content or None if failed
        """
        try:
            file_ext = os.path.splitext(file_path)[1].lower()

            if file_ext in ['.txt', '.md']:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
                if file_ext == '.txt':
                    # Convert plain text to markdown
                    return f"# Document\n\n{content}"
                return content

            elif file_ext == '.csv':
                import pandas as pd
                df = pd.read_csv(file_path, encoding='utf-8', on_bad_lines='skip')
                return f"# CSV Data\n\n{df.to_markdown(index=False)}"

            elif file_ext in ['.html', '.htm']:
                from bs4 import BeautifulSoup
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    soup = BeautifulSoup(f.read(), 'html.parser')
                    text = soup.get_text(separator='\n', strip=True)
                return f"# HTML Document\n\n{text}"

            return None

        except Exception as e:
            logger.error(f"Failed to process text file {file_path}: {e}")
            return None

    async def _upload_to_marker(self, file_path: str, webhook_url: str) -> Optional[str]:
        """
        Upload a file to the Marker API for processing using webhook.

        Args:
            file_path: Path to the file to upload
            webhook_url: URL for webhook callback

        Returns:
            Markdown content or None if failed
        """
        file_ext = os.path.splitext(file_path)[1].lower()
        temp_file_created = False

        try:
            if file_ext == '.pdf':
                mime_type = 'application/pdf'
                process_path = file_path
            elif file_ext in Config.SUPPORTED_IMAGE_FORMATS:
                mime_type = 'application/pdf'
                process_path = self._convert_image_to_pdf(file_path)
                temp_file_created = True
            elif file_ext in Config.SUPPORTED_WORD_FORMATS:
                if file_ext == '.docx':
                    mime_type = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
                else:  # .doc
                    mime_type = 'application/msword'
                process_path = file_path
            elif file_ext in Config.SUPPORTED_EXCEL_FORMATS:
                if file_ext == '.xlsx':
                    mime_type = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
                else:  # .xls
                    mime_type = 'application/vnd.ms-excel'
                process_path = file_path
            else:
                logger.error(f"Unsupported file format: {file_ext}")
                return None

            with open(process_path, 'rb') as f:
                files = {'file': (os.path.basename(process_path), f, mime_type)}

                # Add webhook parameters to the request
                data = {**Config.MARKER_PARAMS}

                response = requests.post(
                    self.api_url,
                    headers=self.headers,
                    files=files,
                    data=data
                )

            if response.status_code != 200:
                logger.error(f"API request failed with status {response.status_code}: {response.text}")
                return None

            # Get the request_id from the response
            result = response.json()
            request_id = result.get('request_id')
            if not request_id:
                logger.error("No request_id returned from Marker API")
                return None

            logger.info(f"File uploaded successfully with request_id: {request_id}")

            # Get webhook manager and create request with the actual request_id
            webhook_mgr = get_webhook_manager()
            future = webhook_mgr.create_request(request_id, webhook_url)

            # Wait for webhook callback
            try:
                webhook_result = await asyncio.wait_for(future, timeout=300)  # 5 minute timeout
                markdown_content = webhook_result.get("markdown_content")

                if not markdown_content:
                    logger.error("Webhook returned empty content")
                    return None

                return markdown_content

            except asyncio.TimeoutError:
                logger.error("Webhook processing timed out")
                return None

        except Exception as e:
            logger.error(f"Upload failed: {e}")
            return None

        finally:
            if temp_file_created and os.path.exists(process_path):
                os.unlink(process_path)
                logger.debug(f"Cleaned up temporary file: {process_path}")
    
    
    def _process_with_local_marker(self, file_path: str) -> Optional[str]:
        """
        Process a file using local Marker installation.

        Args:
            file_path: Path to the file to process

        Returns:
            Markdown content or None if failed
        """
        # Check file extension - Marker has issues with PowerPoint files containing WMF images
        file_ext = os.path.splitext(file_path)[1].lower()
        if file_ext in ['.ppt', '.pptx', '.pps']:
            logger.info(f"Skipping Marker for PowerPoint files due to WMF image issues: {file_ext}")
            return None

        # Check if GPU is enabled
        gpu_enabled = os.getenv("GPU_ENABLED", "false").lower() == "true"
        if not gpu_enabled:
            logger.info("GPU not enabled (GPU_ENABLED=false), skipping local Marker processing")
            return None

        try:
            logger.info("Using marker Python API for GPU-optimized processing")
            
            # Import correct marker modules according to PyPI documentation
            from marker.converters.pdf import PdfConverter
            from marker.models import create_model_dict
            from marker.output import text_from_rendered
            
            # Use preloaded models if available
            try:
                from preload_models import get_preloaded_models
                models = get_preloaded_models()
                if models:
                    logger.info(f"✅ Using {len(models)} preloaded models")
                else:
                    logger.info("Loading models...")
                    models = create_model_dict()
            except ImportError:
                logger.info("Loading models...")
                models = create_model_dict()
            
            # Create optimized configuration for faster processing
            from marker.config.parser import ConfigParser
            
            # Optimized GPU configuration for Tesla T4 (15GB VRAM)
            config = {
                # GPU optimization - maximize parallel processing
                "batch_multiplier": 12,  # Increased for better GPU utilization
                "ocr_batch_size": 32,  # Much larger OCR batch for GPU efficiency
                "layout_batch_size": 8,  # Larger layout batch
                "table_rec_batch_size": 8,  # Larger table batch
                "detection_batch_size": 16,  # Larger detection batch
                
                # OCR optimizations - keep GPU busy
                # "ocr_all_pages": True,  # Process all pages in parallel
                "disable_ocr": False,  # Keep OCR but optimize
                "ocr_error_detection": False,  # Skip for speed
                "detect_language": False,  # Skip language detection
                
                # Processing optimizations - reduce CPU work
                "paginate_output": False,  # Faster without pagination
                "disable_image_extraction": True,  # Skip images
                "skip_table_detection": False,  # Keep tables but optimize
                "disable_math_detection": False,  # Keep math detection
                
                # Parallel processing - maximize GPU usage
                "workers": 8,  # More parallel workers
                "ray_workers": 8,  # Ray parallel processing
                "max_parallel_pages": 4,  # Process multiple pages simultaneously
                
                # GPU memory optimization
                "gpu_memory_fraction": 0.8,  # Use 80% of GPU memory
                "force_gpu": True,  # Force GPU usage
                "cuda_device": 0,  # Use first GPU
                
                # Other optimizations
                # "use_llm": True,  # Keep LLM for quality
                # "langs": ["en"],  # Skip language detection
                # "full_document_analysis": True,  # Process entire document
            }
            
            config_parser = ConfigParser(config)
            
            # Create converter with GPU optimizations and config
            converter = PdfConverter(
                artifact_dict=models,
                config=config_parser.generate_config_dict()
            )
            
            # Force GPU memory optimization
            if torch.cuda.is_available():
                torch.cuda.empty_cache()  # Clear GPU memory
                torch.cuda.set_per_process_memory_fraction(0.8)  # Use 80% of GPU memory
                torch.backends.cudnn.benchmark = True  # Optimize for consistent input sizes
                torch.backends.cuda.matmul.allow_tf32 = True  # Use TensorFloat-32
                torch.backends.cudnn.allow_tf32 = True  # Use TensorFloat-32
            
            # Convert PDF to markdown
            logger.info(f"Converting {file_path} with optimized GPU acceleration...")
            logger.info(f"Settings: batch_multiplier=12, ocr_batch_size=64, workers=8, max_parallel_pages=8")
            start_time = time.perf_counter()
            
            # Run conversion
            rendered = converter(file_path)
            text, _, _ = text_from_rendered(rendered)
            
            processing_time = time.perf_counter() - start_time
            logger.info(f"✅ Conversion completed in {processing_time:.2f}s")
            
            return text
            
        except ImportError as e:
            logger.error(f"Marker Python API not available: {e}")
            logger.error("Please ensure marker-pdf is installed: pip install marker-pdf[gpu]")
            return None
        except Exception as e:
            logger.error(f"Error processing with local Marker: {str(e)}")
            return None


async def convert_to_markdown(
    file_path: str,
    use_local: Optional[bool] = None,
    webhook_url: Optional[str] = None
) -> Optional[str]:
    """
    Convert a PDF or image file to markdown format.

    Args:
        file_path: Path to the input file (PDF or supported image format)
        use_local: Whether to use local Marker (True) or API (False). Defaults to True.
        webhook_url: URL for webhook callback (required if use_local is False)

    Returns:
        Markdown content as string, or None if conversion failed

    Raises:
        FileNotFoundError: If the input file doesn't exist
        ValueError: If the file format is not supported
    """
    # Validate input
    file_path = os.path.expanduser(file_path.strip().strip('"').strip("'"))
    file_path = os.path.abspath(file_path)

    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    file_ext = os.path.splitext(file_path)[1].lower()
    if file_ext not in Config.SUPPORTED_DOCUMENT_FORMATS:
        raise ValueError(f"Unsupported file format: {file_ext}. Supported formats: {Config.SUPPORTED_DOCUMENT_FORMATS}")

    # Process file
    processor = DocumentProcessor()
    logger.info(f"Processing file: {file_path}")

    # Check if it's a text-based file first
    if file_ext in Config.SUPPORTED_TEXT_FORMATS:
        logger.info("Processing text-based file")
        markdown_content = processor._process_text_file(file_path)
        if markdown_content:
            return markdown_content
        else:
            raise ValueError(f"Failed to process text file: {file_path}")

    # Check if it's a .doc file that needs conversion to .docx for Marker processing
    if file_ext == '.doc':
        logger.info("Converting .doc to .docx for Marker GPU processing")

        try:
            import subprocess
            import tempfile

            temp_dir = tempfile.mkdtemp()
            base_name = os.path.splitext(os.path.basename(file_path))[0]
            docx_file = os.path.join(temp_dir, f"{base_name}.docx")

            # Convert .doc to .docx using LibreOffice
            cmd = [
                'libreoffice',
                '--headless',
                '--convert-to', 'docx',
                '--outdir', temp_dir,
                file_path
            ]

            result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)

            if result.returncode == 0 and os.path.exists(docx_file):
                logger.info(f"Successfully converted .doc to .docx: {docx_file}")

                # Update file_path to the .docx file for Marker processing
                file_path = docx_file
                file_ext = '.docx'
                # Continue to Marker processing below

            else:
                logger.error(f"LibreOffice conversion failed: {result.stderr}")
                # Fall back to FastFileExtractor if conversion fails
                logger.info("Falling back to FastFileExtractor for .doc file")
                from integration import FastFileExtractor
                extractor = FastFileExtractor(include_metadata=False)

                result = extractor.extract(file_path)
                if result['success']:
                    return result['markdown']
                else:
                    raise ValueError(f"Both .doc conversion and FastFileExtractor failed")

        except subprocess.TimeoutExpired:
            logger.error("LibreOffice conversion timed out, falling back to FastFileExtractor")
            from integration import FastFileExtractor
            extractor = FastFileExtractor(include_metadata=False)

            result = extractor.extract(file_path)
            if result['success']:
                return result['markdown']
            else:
                raise ValueError(f"LibreOffice timeout and FastFileExtractor failed")

        except Exception as e:
            logger.error(f".doc to .docx conversion failed: {e}, falling back to FastFileExtractor")
            from integration import FastFileExtractor
            extractor = FastFileExtractor(include_metadata=False)

            try:
                result = extractor.extract(file_path)
                if result['success']:
                    return result['markdown']
                else:
                    raise ValueError(f".doc conversion and FastFileExtractor failed: {e}")
            except Exception as fallback_error:
                logger.error(f"FastFileExtractor fallback also failed: {fallback_error}")
                raise ValueError(f"Failed to process .doc file: {e}")

    # Check if it's PowerPoint that should use FastFileExtractor (Marker doesn't support PowerPoint)
    if file_ext in Config.SUPPORTED_POWERPOINT_FORMATS:
        logger.info(f"Processing {file_ext} file with FastFileExtractor (PowerPoint not supported by Marker)")
        from integration import FastFileExtractor
        extractor = FastFileExtractor(include_metadata=False)

        try:
            result = extractor.extract(file_path)
            if result['success']:
                return result['markdown']
            else:
                raise ValueError(f"FastFileExtractor failed: {result.get('error', 'Unknown error')}")
        except Exception as e:
            logger.error(f"FastFileExtractor failed for {file_ext}: {e}")
            raise ValueError(f"Failed to process {file_ext} file: {e}")

    # Default to local processing if use_local is not specified
    if use_local is None:
        use_local = True

    # Choose between local and API processing
    if use_local:
        logger.info("Using local Marker for processing")
        markdown_content = processor._process_with_local_marker(file_path)

        # If local processing fails, raise error instead of fallback
        if markdown_content is None:
            raise RuntimeError("Local Marker processing failed. Please check Marker installation and configuration.")
    else:
        if not webhook_url:
            raise ValueError("webhook_url is required when use_local is False")

        logger.info("Using Marker API with webhook for processing")
        # Upload and wait for webhook response
        markdown_content = await processor._upload_to_marker(file_path, webhook_url)
        if not markdown_content:
            return None

    if not markdown_content:
        return None

    return markdown_content


# new class and functions for fast conversion 
#create another class called Fast_DocumentProcessor
class FastFileExtractor:
    """
    Fast file extraction class for converting various document types to markdown.
    
    Supports PDFs, Word documents, Excel files, PowerPoint presentations, 
    HTML files, CSV files, and other document formats.
    """
    
    FILE_LOADERS = {
        '.pdf': 'pdf', '.doc': 'word', '.docx': 'word', '.xls': 'excel', '.xlsx': 'excel',
        '.ppt': 'powerpoint', '.pptx': 'powerpoint', '.jpg': 'image', '.jpeg': 'image',
        '.png': 'image', '.gif': 'image', '.bmp': 'image', '.tiff': 'image',
        '.html': 'html', '.htm': 'html', '.md': 'markdown', '.csv': 'csv',
        '.epub': 'epub', '.odt': 'odt', '.rtf': 'rtf', '.txt': 'text'
    }
    
    def __init__(self, include_metadata: bool = True):
        """
        Initialize the FastFileExtractor.

        Args:
            include_metadata: Whether to include metadata in results
        """
        self.include_metadata = include_metadata
    
    def extract(self, file_path: str) -> Dict[str, any]:
        """
        Extract content from a file and convert to markdown.
        
        Args:
            file_path: Path to the file to extract
            
        Returns:
            Dictionary containing:
            - markdown: Extracted content as markdown string
            - metadata: File metadata and extraction info
            - processing_time: Time taken for extraction
            - success: Boolean indicating if extraction was successful
            - error: Error message if extraction failed
        """
        start_time = time.time()
        
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        file_ext = Path(file_path).suffix.lower()
        file_type = self.FILE_LOADERS.get(file_ext, 'unknown')
        
        if file_type == 'unknown':
            raise ValueError(f"Unsupported file type: {file_ext}")
        
        try:
            if file_type == 'pdf':
                markdown, metadata = self._extract_pdf(file_path)
            elif file_type == 'word':
                markdown, metadata = self._extract_word(file_path)
            elif file_type == 'excel':
                markdown, metadata = self._extract_excel(file_path)
            elif file_type == 'powerpoint':
                markdown, metadata = self._extract_powerpoint(file_path)
            elif file_type == 'image':
                raise ValueError("Image files not supported. Convert to PDF first.")
            elif file_type in ['html', 'csv', 'markdown', 'text']:
                markdown, metadata = self._extract_simple(file_path, file_type)
            else:
                markdown, metadata = self._extract_unstructured(file_path, file_type)
            
            processing_time = time.time() - start_time
            
            if self.include_metadata:
                metadata.update({
                    'processing_time': f"{processing_time:.2f}s",
                    'file_path': file_path,
                    'file_type': file_type,
                    'file_size': f"{os.path.getsize(file_path) / 1024:.1f} KB"
                })
            
            return {
                'markdown': markdown,
                'metadata': metadata,
                'processing_time': processing_time,
                'success': True
            }
            
        except Exception as e:
            return {
                'markdown': '',
                'metadata': {'error': str(e)},
                'processing_time': time.time() - start_time,
                'success': False,
                'error': str(e)
            }
    
    def _extract_pdf(self, file_path: str) -> Tuple[str, Dict]:
        """Extract content from PDF file using PyMuPDF for fast extraction."""
        try:
            loader = PyMuPDFLoader(file_path)
            docs = loader.load()

            markdown = ""
            for i, doc in enumerate(docs):
                markdown += f"\n{{page_{i}}}\n---\n\n{doc.page_content}\n\n"

            return markdown, {'total_pages': len(docs), 'loader': 'PyMuPDF'}
        except Exception as e:
            logger.warning(f"PyMuPDFLoader failed: {e}")
            # Fallback to Unstructured if PyMuPDF fails
            try:
                loader = UnstructuredPDFLoader(
                    file_path,
                    mode="single",
                    strategy="fast"
                )
                docs = loader.load()
                markdown = self._format_documents(docs)
                return markdown, {'total_elements': len(docs), 'loader': 'Unstructured'}
            except Exception as e2:
                logger.error(f"UnstructuredPDFLoader also failed: {e2}")
                # Last resort: Try basic text extraction
                try:
                    import PyPDF2
                    markdown = "# PDF Document\n\n"
                    with open(file_path, 'rb') as file:
                        reader = PyPDF2.PdfReader(file)
                        for page_num, page in enumerate(reader.pages):
                            text = page.extract_text()
                            if text:
                                markdown += f"\n## Page {page_num + 1}\n\n{text}\n\n"
                    return markdown, {'total_pages': len(reader.pages), 'loader': 'PyPDF2'}
                except Exception as e3:
                    logger.error(f"PyPDF2 also failed: {e3}")
                    raise e2
    
    def _extract_word(self, file_path: str) -> Tuple[str, Dict]:
        """Extract content from Word document with multiple fallbacks."""
        # First, try using local processing without API calls
        try:
            loader = UnstructuredWordDocumentLoader(
                file_path,
                mode="single",
                strategy="fast",  # Use fast local strategy
                hi_res_model_name=None,  # Disable high-res model that might use API
                use_api=False  # Explicitly disable API usage if supported
            )
            docs = loader.load()
            return self._format_documents(docs), {'total_elements': len(docs)}
        except Exception as e:
            logger.warning(f"UnstructuredWordDocumentLoader failed: {e}")

        # Fallback 1: Try using python-docx directly for .docx files
        if file_path.endswith('.docx'):
            try:
                import docx
                doc = docx.Document(file_path)
                markdown = "# Word Document\n\n"
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        markdown += paragraph.text + "\n\n"

                # Also extract tables
                for table in doc.tables:
                    markdown += "\n"
                    for row in table.rows:
                        markdown += "| " + " | ".join([cell.text.strip() for cell in row.cells]) + " |\n"
                        if table.rows.index(row) == 0:
                            markdown += "|" + "---|" * len(row.cells) + "\n"
                    markdown += "\n"

                return markdown, {'total_elements': len(doc.paragraphs), 'loader': 'python-docx'}
            except Exception as e2:
                logger.error(f"python-docx fallback failed: {e2}")

        # Fallback 2: Try extracting raw XML text for corrupted .docx
        if file_path.endswith('.docx'):
            try:
                import zipfile
                import xml.etree.ElementTree as ET

                markdown = "# Word Document (Text Recovery)\n\n"
                with zipfile.ZipFile(file_path, 'r') as docx:
                    # Extract text from document.xml
                    with docx.open('word/document.xml') as xml_file:
                        tree = ET.parse(xml_file)
                        root = tree.getroot()

                        # Extract all text elements
                        namespace = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
                        paragraphs = root.findall('.//w:t', namespace)

                        current_paragraph = []
                        for elem in root.findall('.//w:p', namespace):
                            texts = elem.findall('.//w:t', namespace)
                            para_text = ' '.join([t.text for t in texts if t.text])
                            if para_text.strip():
                                markdown += para_text.strip() + "\n\n"

                return markdown, {'total_elements': 1, 'loader': 'xml-recovery'}
            except Exception as e3:
                logger.error(f"XML extraction fallback failed: {e3}")

        # Fallback 3: For .doc files, try using antiword or other text extraction
        if file_path.endswith('.doc'):
            try:
                import subprocess
                result = subprocess.run(['antiword', file_path], capture_output=True, text=True, timeout=30)
                if result.returncode == 0 and result.stdout:
                    markdown = "# Word Document\n\n" + result.stdout
                    return markdown, {'total_elements': 1, 'loader': 'antiword'}
            except Exception as e4:
                logger.error(f"antiword fallback failed: {e4}")

        # Last resort: Try to extract any readable text
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
                # Clean up binary garbage
                import string
                printable = set(string.printable)
                content = ''.join(filter(lambda x: x in printable, content))
                if len(content) > 100:  # Only if we got meaningful content
                    markdown = "# Document (Text Recovery)\n\n" + content
                    return markdown, {'total_elements': 1, 'loader': 'text-recovery'}
        except Exception as e5:
            logger.error(f"Text recovery failed: {e5}")

        # If all fallbacks fail, raise error
        raise ValueError(f"Failed to extract Word document after all fallback attempts")
    
    def _extract_excel(self, file_path: str) -> Tuple[str, Dict]:
        """Extract content from Excel file with robust error handling."""
        logger.info(f"Processing Excel file: {file_path}")

        # Try multiple approaches to handle corrupted Excel files

        # Approach 1: Try pandas with different engines
        try:
            import pandas as pd

            # Try reading with openpyxl engine first
            try:
                xls = pd.ExcelFile(file_path, engine='openpyxl')
            except:
                # If openpyxl fails, try xlrd for older formats
                try:
                    xls = pd.ExcelFile(file_path, engine='xlrd')
                except:
                    # Last resort: try calamine engine (if available)
                    try:
                        xls = pd.ExcelFile(file_path, engine='calamine')
                    except:
                        raise

            markdown = "# Excel Document\n\n"

            for sheet_name in xls.sheet_names:
                try:
                    df = pd.read_excel(xls, sheet_name=sheet_name)
                    markdown += f"## Sheet: {sheet_name}\n\n"

                    if not df.empty:
                        # Limit columns and rows for very large sheets
                        if len(df) > 1000:
                            df = df.head(1000)
                            markdown += "*Note: Showing first 1000 rows*\n\n"
                        if len(df.columns) > 50:
                            df = df.iloc[:, :50]
                            markdown += "*Note: Showing first 50 columns*\n\n"

                        markdown += df.to_markdown(index=False) + "\n\n"
                    else:
                        markdown += "*(Empty sheet)*\n\n"
                except Exception as sheet_error:
                    markdown += f"## Sheet: {sheet_name}\n\n*(Error reading sheet: {str(sheet_error)})*\n\n"
                    logger.warning(f"Failed to read sheet {sheet_name}: {sheet_error}")

            return markdown, {'total_sheets': len(xls.sheet_names), 'loader': 'pandas'}

        except ImportError:
            logger.error("pandas not installed for Excel processing")
        except Exception as e:
            logger.error(f"pandas failed: {e}")

            # Last resort: try openpyxl directly for .xlsx files
            if file_path.endswith('.xlsx'):
                try:
                    from openpyxl import load_workbook
                    wb = load_workbook(file_path, read_only=True, data_only=True)
                    markdown = "# Excel Document\n\n"

                    for sheet_name in wb.sheetnames:
                        sheet = wb[sheet_name]
                        markdown += f"## Sheet: {sheet_name}\n\n"

                        # Extract data from sheet
                        data = []
                        for row in sheet.iter_rows(values_only=True):
                            if any(cell is not None for cell in row):
                                data.append([str(cell) if cell is not None else "" for cell in row])

                        if data:
                            # Format as table
                            markdown += "| " + " | ".join(data[0]) + " |\n"
                            markdown += "|" + "---|" * len(data[0]) + "\n"
                            for row in data[1:]:
                                markdown += "| " + " | ".join(row) + " |\n"
                            markdown += "\n"
                        else:
                            markdown += "*(Empty sheet)*\n\n"

                    wb.close()
                    return markdown, {'total_sheets': len(wb.sheetnames), 'loader': 'openpyxl'}
                except ImportError:
                    logger.error("openpyxl not installed for Excel fallback")
                except Exception as e2:
                    logger.error(f"openpyxl fallback failed: {e2}")

        # Fallback 2: Try converting Excel to CSV using LibreOffice, then process as CSV
        try:
            logger.info("Attempting Excel to CSV conversion using LibreOffice...")
            import subprocess
            import tempfile

            # Create temp directory for conversion
            temp_dir = tempfile.mkdtemp()
            base_name = os.path.splitext(os.path.basename(file_path))[0]

            # Convert Excel to CSV using LibreOffice
            cmd = [
                'libreoffice',
                '--headless',
                '--convert-to', 'csv',
                '--outdir', temp_dir,
                file_path
            ]

            result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)

            if result.returncode == 0:
                # Look for the converted CSV file
                csv_file = os.path.join(temp_dir, f"{base_name}.csv")
                if os.path.exists(csv_file):
                    logger.info(f"Successfully converted Excel to CSV: {csv_file}")

                    # Process the CSV file
                    try:
                        import pandas as pd
                        df = pd.read_csv(csv_file, encoding='utf-8', on_bad_lines='skip')
                        markdown = "# Excel Document (Converted from CSV)\n\n"

                        if not df.empty:
                            # Limit very large datasets
                            if len(df) > 1000:
                                df = df.head(1000)
                                markdown += "*Note: Showing first 1000 rows*\n\n"
                            if len(df.columns) > 50:
                                df = df.iloc[:, :50]
                                markdown += "*Note: Showing first 50 columns*\n\n"

                            markdown += df.to_markdown(index=False) + "\n\n"
                        else:
                            markdown += "*(No data found)*\n\n"

                        # Clean up temp files
                        try:
                            os.unlink(csv_file)
                            os.rmdir(temp_dir)
                        except:
                            pass

                        return markdown, {'total_sheets': 1, 'loader': 'libreoffice-csv'}

                    except Exception as csv_error:
                        logger.error(f"Failed to process converted CSV: {csv_error}")
                        # Clean up and continue to next fallback
                        try:
                            os.unlink(csv_file)
                            os.rmdir(temp_dir)
                        except:
                            pass
            else:
                logger.error(f"LibreOffice conversion failed: {result.stderr}")

        except subprocess.TimeoutExpired:
            logger.error("LibreOffice conversion timed out")
        except Exception as libre_error:
            logger.error(f"LibreOffice conversion error: {libre_error}")

        # Fallback 3: Try using xlwings for complex Excel files (Windows only)
        if os.name == 'nt':  # Windows only
            try:
                logger.info("Attempting xlwings extraction (Windows only)...")
                import xlwings as xw

                app = xw.App(visible=False)
                wb = app.books.open(file_path)
                markdown = "# Excel Document (xlwings extraction)\n\n"

                for sheet in wb.sheets:
                    try:
                        markdown += f"## Sheet: {sheet.name}\n\n"

                        # Get used range
                        used_range = sheet.used_range
                        if used_range:
                            # Convert to pandas DataFrame for easier handling
                            values = used_range.value
                            if values:
                                import pandas as pd
                                if isinstance(values[0], list):
                                    df = pd.DataFrame(values)
                                else:
                                    df = pd.DataFrame([values])

                                # Limit size
                                if len(df) > 1000:
                                    df = df.head(1000)
                                    markdown += "*Note: Showing first 1000 rows*\n\n"

                                markdown += df.to_markdown(index=False) + "\n\n"
                            else:
                                markdown += "*(Empty sheet)*\n\n"
                        else:
                            markdown += "*(Empty sheet)*\n\n"

                    except Exception as sheet_error:
                        markdown += f"*(Error reading sheet {sheet.name}: {str(sheet_error)})*\n\n"
                        logger.warning(f"xlwings failed to read sheet {sheet.name}: {sheet_error}")

                wb.close()
                app.quit()

                return markdown, {'total_sheets': len(wb.sheets), 'loader': 'xlwings'}

            except ImportError:
                logger.warning("xlwings not available (Windows only or not installed)")
            except Exception as xlwings_error:
                logger.error(f"xlwings extraction failed: {xlwings_error}")
                try:
                    # Ensure Excel app is closed
                    if 'app' in locals():
                        app.quit()
                except:
                    pass

        # Fallback 4: Raw text extraction as last resort
        try:
            logger.info("Attempting raw text extraction from Excel file...")
            # For .xlsx files, try to extract text from the zip archive
            if file_path.endswith('.xlsx'):
                import zipfile
                import xml.etree.ElementTree as ET

                markdown = "# Excel Document (Raw Text Extraction)\n\n"
                with zipfile.ZipFile(file_path, 'r') as xlsx:
                    # Try to extract from shared strings and worksheets
                    try:
                        with xlsx.open('xl/sharedStrings.xml') as shared_strings:
                            tree = ET.parse(shared_strings)
                            root = tree.getroot()

                            # Extract text from shared strings
                            strings = []
                            for si in root.findall('.//{http://schemas.openxmlformats.org/spreadsheetml/2006/main}t'):
                                if si.text:
                                    strings.append(si.text)

                            if strings:
                                markdown += "**Extracted Text Content:**\n\n"
                                markdown += " | ".join(strings[:100])  # Limit to first 100 strings
                                if len(strings) > 100:
                                    markdown += f"\n\n*(... and {len(strings) - 100} more entries)*"
                                markdown += "\n\n"

                    except:
                        pass

                    # Also try to extract from worksheet files
                    worksheet_files = [f for f in xlsx.namelist() if f.startswith('xl/worksheets/')]
                    for ws_file in worksheet_files[:3]:  # Limit to first 3 worksheets
                        try:
                            with xlsx.open(ws_file) as worksheet:
                                tree = ET.parse(worksheet)
                                root = tree.getroot()

                                # Extract cell values
                                values = []
                                for v in root.findall('.//{http://schemas.openxmlformats.org/spreadsheetml/2006/main}v'):
                                    if v.text:
                                        values.append(v.text)

                                if values:
                                    markdown += f"**Worksheet {ws_file}:**\n\n"
                                    markdown += " | ".join(values[:50])  # Limit to first 50 values
                                    markdown += "\n\n"
                        except:
                            continue

                if len(markdown) > 50:  # If we extracted some content
                    return markdown, {'total_sheets': len(worksheet_files), 'loader': 'raw-xml'}

        except Exception as raw_error:
            logger.error(f"Raw text extraction failed: {raw_error}")

        # If all fallbacks fail, raise the original error
        raise ValueError(f"Failed to extract Excel file after all fallback attempts: {e}")
    
    def _extract_powerpoint(self, file_path: str) -> Tuple[str, Dict]:
        """Extract content from PowerPoint presentation using python-pptx directly."""
        logger.info(f"Processing PowerPoint file: {file_path}")

        # Use python-pptx directly to avoid API calls
        if file_path.endswith('.pptx'):
            try:
                from pptx import Presentation

                prs = Presentation(file_path)
                markdown = ""

                for slide_num, slide in enumerate(prs.slides, 1):
                    # Extract text from all shapes
                    slide_content = ""
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_content += shape.text.strip() + "\n\n"

                        # Handle tables
                        if shape.has_table:
                            table = shape.table
                            for row_idx, row in enumerate(table.rows):
                                slide_content += "| " + " | ".join([cell.text for cell in row.cells]) + " |\n"
                                if row_idx == 0:
                                    slide_content += "|" + "---|" * len(row.cells) + "\n"
                            slide_content += "\n"

                    # Only add slide content if it has text
                    if slide_content.strip():
                        markdown += slide_content

                return markdown, {'total_slides': len(prs.slides), 'loader': 'python-pptx'}

            except Exception as e:
                logger.error(f"python-pptx failed: {e}")

                # Fallback to UnstructuredPowerPointLoader as last resort
                try:
                    loader = UnstructuredPowerPointLoader(
                        file_path,
                        mode="single",
                        strategy="fast"
                    )
                    docs = loader.load()

                    markdown = ""
                    for doc in docs:
                        markdown += doc.page_content + "\n\n"

                    return markdown, {'total_elements': len(docs), 'loader': 'UnstructuredPowerPointLoader'}

                except Exception as e2:
                    logger.error(f"UnstructuredPowerPointLoader also failed: {e2}")
                    raise ValueError(f"Failed to extract PowerPoint file: {e}")

        # For .ppt files (older format), try UnstructuredPowerPointLoader
        else:
            try:
                loader = UnstructuredPowerPointLoader(
                    file_path,
                    mode="single",
                    strategy="fast"
                )
                docs = loader.load()

                markdown = ""
                for doc in docs:
                    markdown += doc.page_content + "\n\n"

                return markdown, {'total_elements': len(docs), 'loader': 'UnstructuredPowerPointLoader'}

            except Exception as e:
                logger.error(f"Failed to extract .ppt file: {e}")
                raise ValueError(f"Failed to extract PowerPoint file: {e}")
    
    def _extract_simple(self, file_path: str, file_type: str) -> Tuple[str, Dict]:
        """Extract content from simple file types (text, markdown, HTML, CSV)."""
        if file_type == 'markdown':
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read(), {'loader': 'Direct'}
        
        elif file_type == 'text':
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            return f"# Text Document\n\n{content}", {'loader': 'Direct'}
        
        elif file_type == 'html':
            try:
                loader = UnstructuredHTMLLoader(
                    file_path,
                    mode="single"
                )
                docs = loader.load()
                return self._format_documents(docs), {'loader': 'HTML'}
            except Exception as e:
                logger.warning(f"UnstructuredHTMLLoader failed: {e}")
                # Fallback: Read HTML directly
                from bs4 import BeautifulSoup
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    soup = BeautifulSoup(f.read(), 'html.parser')
                    # Extract text
                    text = soup.get_text(separator='\n', strip=True)
                    return f"# HTML Document\n\n{text}", {'loader': 'BeautifulSoup'}
        
        elif file_type == 'csv':
            try:
                loader = UnstructuredCSVLoader(
                    file_path,
                    mode="single"
                )
                docs = loader.load()
            except Exception as e:
                logger.warning(f"UnstructuredCSVLoader failed: {e}")
                # Fallback: Use pandas
                try:
                    import pandas as pd
                    df = pd.read_csv(file_path)
                    markdown = "# CSV Data\n\n"
                    markdown += df.to_markdown(index=False) if not df.empty else "*(Empty CSV)*"
                    return markdown, {'loader': 'pandas'}
                except Exception as e2:
                    # Last resort: Read as text
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        content = f.read()
                    return self._format_table(content), {'loader': 'text'}
            markdown = "# CSV Data\n\n"
            for doc in docs:
                markdown += self._format_table(doc.page_content)
            return markdown, {'loader': 'CSV'}
    
    def _extract_unstructured(self, file_path: str, file_type: str) -> Tuple[str, Dict]:
        """Extract content using unstructured loaders for specialized formats."""
        loaders = {
            'epub': UnstructuredEPubLoader,
            'odt': UnstructuredODTLoader,
            'rtf': UnstructuredRTFLoader
        }
        
        loader_class = loaders.get(file_type)
        if not loader_class:
            raise ValueError(f"No loader available for {file_type}")
        
        # Try with minimal parameters
        try:
            loader = loader_class(
                file_path,
                mode="single"
            )
        except TypeError:
            # Some loaders might not accept mode parameter
            loader = loader_class(file_path)
        docs = loader.load()
        return self._format_documents(docs), {'loader': file_type.upper()}
    
    def _format_documents(self, docs: List) -> str:
        """Format document list into markdown string."""
        return "\n\n".join(doc.page_content for doc in docs) + "\n"
    
    def _format_table(self, content: str) -> str:
        """Format tabular content as markdown table."""
        lines = content.strip().split('\n')
        if not lines:
            return content
        
        markdown = "\n"
        for i, line in enumerate(lines):
            if '\t' in line:
                cells = line.split('\t')
                markdown += '| ' + ' | '.join(cells) + ' |\n'
                if i == 0:
                    markdown += '|' + '---|' * len(cells) + '\n'
            else:
                markdown += line + '\n'
        
        return markdown + "\n"
    
    def extract_batch(self, file_paths: List[str]) -> List[Dict]:
        """
        Extract content from multiple files in batch.
        
        Args:
            file_paths: List of file paths to extract
            
        Returns:
            List of extraction results for each file
        """
        results = []
        logger.info(f"Starting batch extraction of {len(file_paths)} files")
        
        for i, file_path in enumerate(file_paths):
            try:
                logger.info(f"Processing file {i+1}/{len(file_paths)}: {Path(file_path).name}")
                result = self.extract(file_path)
                results.append(result)
                
            except Exception as e:
                error_result = {
                    'markdown': '',
                    'metadata': {'error': str(e), 'file_path': file_path},
                    'processing_time': 0,
                    'success': False,
                    'error': str(e)
                }
                results.append(error_result)
                logger.error(f"Failed to extract {file_path}: {e}")
        
        successful = sum(1 for r in results if r['success'])
        logger.info(f"Batch extraction completed: {successful}/{len(file_paths)} files successful")
        return results

def fast_convert_to_markdown(
    file_path: str,
    include_metadata: Optional[bool] = False
) -> Optional[str]:
    """
    Fast convert a document file to markdown format using FastFileExtractor.

    Args:
        file_path: Path to the input file (supports PDF, Word, Excel, PowerPoint, images, etc.)
        include_metadata: Whether to include processing metadata in output
        
    Returns:
        Markdown content as string, or None if conversion failed
        
    Raises:
        FileNotFoundError: If the input file doesn't exist
        ValueError: If the file format is not supported
    """
    # Validate input
    file_path = os.path.expanduser(file_path.strip().strip('"').strip("'"))
    file_path = os.path.abspath(file_path)
    
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")
    
    file_ext = os.path.splitext(file_path)[1].lower()
    if file_ext not in FastFileExtractor.FILE_LOADERS:
        supported_formats = list(FastFileExtractor.FILE_LOADERS.keys())
        raise ValueError(f"Unsupported file format: {file_ext}. Supported formats: {supported_formats}")
    
    # Initialize fast extractor
    extractor = FastFileExtractor(
        include_metadata=include_metadata
    )

    logger.info(f"Fast processing file: {file_path}")
    
    # Extract content
    try:
        result = extractor.extract(file_path)
        
        if not result['success']:
            logger.error(f"Fast conversion failed: {result.get('error', 'Unknown error')}")
            return None
        
        markdown_content = result['markdown']
        processing_time = result['processing_time']
        
        if not markdown_content:
            logger.error("Fast conversion returned empty content")
            return None
        
        logger.info(f"Fast conversion completed in {processing_time:.2f}s")
        
    except Exception as e:
        logger.error(f"Fast conversion error: {e}")
        return None

    return markdown_content

def chunk_markdown(
    markdown_input: Union[str, Path],
) -> List[Dict]:
    """
    Chunk markdown content into smaller pieces for vector storage.
    
    Args:
        markdown_input: Either markdown text string or path to markdown file
        max_chunk_size: Maximum characters per chunk (default: 2500)
        new_chunk_after: Start new chunk after this many characters (default: 2000)
        min_chunk_size: Minimum chunk size to avoid tiny chunks (default: 500)
        chunk_overlap: Character overlap between chunks (default: 200)
        include_metadata: Whether to include metadata in chunk output
        
    Returns:
        List of chunk dictionaries containing content and metadata
        
    Raises:
        FileNotFoundError: If markdown_input is a path that doesn't exist
        ValueError: If the input is empty
    """
    # Handle input - either text or file path
    markdown_content = markdown_input
    if isinstance(markdown_input, Path):
        with open(markdown_input, 'r', encoding='utf-8') as f:
            markdown_content = f.read()
        logger.info(f"Loaded markdown from file: {markdown_input}")
    
    if not markdown_content or not markdown_content.strip():
        raise ValueError("Empty markdown content")
    
    # Set chunking parameters
    max_chunk_size = Config.MAX_CHUNK_SIZE
    new_chunk_after = Config.NEW_CHUNK_AFTER
    min_chunk_size = Config.MIN_CHUNK_SIZE
    chunk_overlap = Config.CHUNK_OVERLAP
    
    # Parse markdown into elements
    try:
        elements = partition_md(text=markdown_content)
    except Exception as e:
        logger.error(f"Failed to partition markdown: {e}")
        raise
    
    # Extract titles for context
    title_by_id = {}
    for el in elements:
        if getattr(el, "category", None) == "Title" and hasattr(el, "element_id"):
            title_by_id[getattr(el, "element_id")] = str(el).strip()
    
    # Chunk the elements
    try:
        if chunk_by_title is None:
            raise ImportError("chunk_by_title not available")
        chunks = chunk_by_title(
            elements,
            max_characters=max_chunk_size,
            new_after_n_chars=new_chunk_after,
            combine_text_under_n_chars=min_chunk_size
        )
        logger.info(f"Used title-based chunking")
    except Exception as e:
        logger.warning(f"Title-based chunking failed, falling back to basic: {e}")
        chunks = chunk_elements(
            elements,
            max_characters=max_chunk_size,
            new_after_n_chars=new_chunk_after,
            overlap=chunk_overlap
        )
    
    # Prepare final chunks with metadata
    tokenizer = tiktoken.encoding_for_model("gpt-3.5-turbo")
    final_chunks = []
    
    for i, chunk in enumerate(chunks):
        text = str(chunk).strip()
        if not text:
            continue
        
        # Add header context if available
        header = None
        
        # Use Unstructured's parent_id to retrieve title directly
        parent_id = getattr(getattr(chunk, "metadata", None), "parent_id", None)
        if parent_id and parent_id in title_by_id:
            header = title_by_id[parent_id][:200]
        
        # If no parent relationship found, check if chunk starts with a title-like line
        if not header and hasattr(chunk, 'category') and chunk.category == 'CompositeElement':
            first_line = text.split('\n')[0].strip()
            # Simple check: if first line is short and looks like a title
            if len(first_line) < 100 and len(first_line) > 5:
                # Check if it's mostly uppercase or title case
                if (first_line.isupper() or 
                    first_line.istitle() or 
                    first_line.count(':') > 0):  # Often titles end with colons
                    header = first_line
        
        # For TableChunk, use a descriptive identifier
        if hasattr(chunk, 'category') and chunk.category == 'TableChunk':
            header = "Table Data"
        
        chunk_data = {
            "content": text,
            "chunk_index": i,
            "header": header,
            "char_count": len(text),
            "token_count": len(tokenizer.encode(text)),
            "chunk_type": chunk.category if hasattr(chunk, 'category') else "text"
        }
        
        final_chunks.append(chunk_data)
    
    logger.info(f"Created {len(final_chunks)} chunks from markdown content")
    
    return final_chunks

def save_document_to_weaviate(
    client,
    file_path: str,
    chunks: List[Dict],
    document_id: str,
    tenant_id: str
) -> str:
    """
    Save document metadata to Weaviate Documents collection.
    
    Args:
        client: Weaviate client instance
        file_path: Path to the original document file
        chunks: List of chunks (used to get total_chunks count)
        document_id: Unique identifier for the document
        tenant_id: Tenant ID for multi-tenancy
        
    Returns:
        The document UUID used in Weaviate
    """
    # Get file information
    file_stats = os.stat(file_path)
    file_path_obj = Path(file_path)
    
    file_name = file_path_obj.name
    file_size = file_stats.st_size
    mime_type = mimetypes.guess_type(file_path)[0]
    total_chunks = len(chunks)
    
    # Prepare document data
    document_data = {
        "file_name": file_name,
        "file_size": file_size,
        "total_chunks": total_chunks,
        "mime_type": mime_type or "application/octet-stream",
    }
    
    # Save to Weaviate Documents collection
    documents_collection = client.collections.get("Documents").with_tenant(tenant_id)
    documents_collection.data.insert(
        properties=document_data,
        uuid=document_id
    )
    
    logger.info(f"Saved document '{file_name}' to Weaviate with UUID: {document_id}")
    return document_id

def save_chunks_to_weaviate(
    client,
    chunks: List[Dict],
    document_id: str,
    tenant_id: str
) -> None:
    """
    Save chunks to Weaviate Chunks collection with reference to document.
    
    Args:
        client: Weaviate client instance
        chunks: List of chunk dictionaries from chunk_markdown
        document_id: Document UUID to reference
        tenant_id: Tenant ID for multi-tenancy
    """
    chunks_collection = client.collections.get("Chunks").with_tenant(tenant_id)
    
    with chunks_collection.batch.fixed_size(batch_size=50) as batch:
        for i, chunk in enumerate(chunks):
            chunk_uuid = generate_uuid5(f"{document_id}_chunk_{i}")
            
            chunk_data = {
                "content": chunk["content"],
                "chunk_index": i,
                "header": chunk.get("header", ""),
                "document_id": document_id,
                "char_count": chunk.get("char_count", len(chunk["content"])),
                "token_count": chunk.get("token_count", 0),
                "chunk_type": chunk.get("chunk_type", "text")
            }
            
            batch.add_object(
                properties=chunk_data,
                uuid=chunk_uuid,
                references={"document_object": document_id}
            )
    
    logger.info(f"Saved {len(chunks)} chunks to Weaviate for document {document_id}")

def delete_document_from_weaviate(
    document_id: str,
    tenant_id: str,
    client: Optional[object] = None
) -> Dict[str, Union[str, int]]:
    """
    Delete a document and all its associated chunks from Weaviate.
    
    Args:
        document_id: UUID of the document to delete
        tenant_id: Tenant ID for multi-tenancy
        client: Weaviate client instance (optional, will create if not provided)
        
    Returns:
        Dictionary containing:
        - document_id: The document UUID that was deleted
        - chunks_deleted: Number of chunks deleted
        - status: Success message
        
    Raises:
        Exception: If deletion fails
    """
    # Create client if not provided
    client_created = False
    if client is None:
        client = _get_weaviate_client()
        client_created = True
    
    try:
        # Track deletion results
        chunks_deleted = 0
        
        # Step 1: Delete all chunks associated with this document
        chunks_collection = client.collections.get("Chunks").with_tenant(tenant_id)
        
        # Use where filter to find and delete all chunks with this document_id
        logger.info(f"Deleting chunks for document {document_id}")
        
        # First, query to find chunks with this document_id
        try:
            # Use the correct query syntax for newer Weaviate
            from weaviate.classes.query import Filter
            
            chunk_results = chunks_collection.query.fetch_objects(
                where=Filter.by_property("document_id").equal(document_id),
                limit=1000  # Adjust if documents have more chunks
            )
            
            # Delete each chunk by UUID
            for chunk in chunk_results.objects:
                chunks_collection.data.delete_by_id(chunk.uuid)
                chunks_deleted += 1
                
            logger.info(f"Deleted {chunks_deleted} chunks")
        except Exception as e:
            logger.warning(f"Could not query/delete chunks: {e}")
            # Try alternative deletion method
            try:
                # Delete by batch using where condition
                deleted = chunks_collection.data.delete_many(
                    where=Filter.by_property("document_id").equal(document_id)
                )
                chunks_deleted = deleted.successful
                logger.info(f"Batch deleted {chunks_deleted} chunks")
            except Exception as e2:
                logger.warning(f"Batch deletion also failed: {e2}")
                # Continue to delete document even if chunks deletion has issues
        
        # Step 2: Delete the document itself
        documents_collection = client.collections.get("Documents").with_tenant(tenant_id)
        
        try:
            documents_collection.data.delete_by_id(document_id)
            logger.info(f"Deleted document {document_id}")
        except Exception as e:
            logger.error(f"Failed to delete document: {e}")
            raise
        
        result = {
            "document_id": document_id,
            "chunks_deleted": chunks_deleted,
            "status": f"Successfully deleted document and {chunks_deleted} chunks"
        }
        
        logger.info(f"Successfully deleted document {document_id} and {chunks_deleted} associated chunks")
        return result
        
    except Exception as e:
        logger.error(f"Error deleting document: {e}")
        raise
    finally:
        # Close client if we created it
        if client_created:
            try:
                client.close()
                logger.info("Closed Weaviate connection")
            except Exception as e:
                logger.warning(f"Error closing Weaviate connection: {e}")

def _get_weaviate_client(
    api_key: Optional[str] = None,
    url: Optional[str] = None,
    http_port: Optional[int] = None,
    grpc_port: Optional[int] = None
):
    """
    Create and connect to a Weaviate client instance.
    
    Args:
        api_key: Weaviate API key (defaults to env variable)
        url: Weaviate URL (defaults to env variable)
        http_port: HTTP port (defaults to 8080)
        grpc_port: GRPC port (defaults to 50051)
        
    Returns:
        Connected WeaviateClient instance
        
    Raises:
        Exception: If connection fails
    """
    if WeaviateClient is None:
        logger.warning("WeaviateClient is not available. Please install weaviate-client package.")
        return None
    
    api_key = api_key or Config.WEAVIATE_API_KEY
    url = url or Config.WEAVIATE_URL
    http_port = http_port or Config.WEAVIATE_HTTP_PORT
    grpc_port = grpc_port or Config.WEAVIATE_GRPC_PORT
    
    try:
        # Try newer weaviate-client API
        client = WeaviateClient(
            connection_params=ConnectionParams(
                http=ProtocolParams(host=url, port=http_port, secure=False),
                grpc=ProtocolParams(host=url, port=grpc_port, secure=False)
            ),
            auth_client_secret=AuthApiKey(api_key),
        )
        
        client.connect()
        logger.info(f"Connected to Weaviate at {url}")
        return client
    except Exception:
        # Fallback to older weaviate-client API
        try:
            import weaviate
            client = weaviate.Client(
                url=f"http://{url}:{http_port}",
                auth_client_secret=weaviate.AuthApiKey(api_key)
            )
            logger.info(f"Connected to Weaviate at {url} (legacy client)")
            return client
        except Exception as e2:
            logger.error(f"Failed to connect to Weaviate: {e2}")
            raise

async def process_document_to_weaviate(
    file_path: str,
    document_id: str,
    tenant_id: str,
    use_local: Optional[bool] = True,
    webhook_url: Optional[str] = None,
    client: Optional[object] = None
) -> Dict[str, Union[str, int, List[Dict]]]:
    """
    Complete pipeline to process a document and save to Weaviate.
    
    This function:
    1. Converts the document to markdown
    2. Chunks the markdown content
    3. Saves document metadata to Weaviate
    4. Saves chunks to Weaviate with references

    Args:
        file_path: Path to the input document (PDF or image)
        document_id: Unique identifier for the document
        tenant_id: Tenant ID for multi-tenancy
        use_local: Whether to use local Marker (True) or API (False). Defaults to True.
        client: Weaviate client instance (optional, will create if not provided)
        
    Returns:
        Dictionary containing:
        - document_id: The document UUID used
        - file_name: Original file name
        - total_chunks: Number of chunks created
        - chunks: List of chunk dictionaries (if needed for further processing)
        
    Raises:
        FileNotFoundError: If the input file doesn't exist
        ValueError: If conversion or chunking fails
        Exception: If Weaviate operations fail
    """
    # Create client if not provided
    client_created = False
    if client is None:
        client = _get_weaviate_client()
        client_created = True
    
    try:
        # Step 1: Convert document to markdown
        logger.info(f"Processing document: {file_path}")

        markdown_content = await convert_to_markdown(
            file_path=file_path,
            use_local=use_local,
            webhook_url=webhook_url
        )

        if not markdown_content:
            raise ValueError(f"Failed to convert document to markdown: {file_path}")

        # Step 2: Chunk the markdown content
        chunks = chunk_markdown(
            markdown_input=markdown_content
        )

        if not chunks:
            raise ValueError(f"No chunks created from document: {file_path}")
        
        # Step 3: Save document metadata to Weaviate
        saved_doc_id = save_document_to_weaviate(
            client=client,
            file_path=file_path,
            chunks=chunks,
            document_id=document_id,
            tenant_id=tenant_id
        )
        
        # Step 4: Save chunks to Weaviate
        save_chunks_to_weaviate(
            client=client,
            chunks=chunks,
            document_id=saved_doc_id,
            tenant_id=tenant_id
        )
        
        # Return summary information
        result = {
            "document_id": saved_doc_id,
            "file_name": Path(file_path).name,
            "total_chunks": len(chunks),
            "chunks": chunks  # Include chunks in case caller needs them
        }
        
        logger.info(f"Successfully processed document '{Path(file_path).name}' with {len(chunks)} chunks")
        return result
        
    except FileNotFoundError as e:
        logger.error(f"File not found: {e}")
        raise
    except ValueError as e:
        logger.error(f"Processing error: {e}")
        raise
    except Exception as e:
        logger.error(f"Unexpected error processing document: {e}")
        raise
    finally:
        # Close client if we created it
        if client_created:
            try:
                client.close()
                logger.info("Closed Weaviate connection")
            except Exception as e:
                logger.warning(f"Error closing Weaviate connection: {e}")


def fast_doc_to_weaviate(
    file_path: str,
    document_id: str,
    tenant_id: str,
    client: Optional[object] = None,
    include_metadata: Optional[bool] = False
) -> Dict[str, Union[str, int, List[Dict]]]:
    """
    Fast pipeline to process a document and save to Weaviate using fast_convert_to_markdown.
    
    This function:
    1. Converts the document to markdown using fast_convert_to_markdown
    2. Chunks the markdown content
    3. Saves document metadata to Weaviate
    4. Saves chunks to Weaviate with references
    
    Args:
        file_path: Path to the input document (PDF or image)
        document_id: Unique identifier for the document
        tenant_id: Tenant ID for multi-tenancy
        client: Weaviate client instance (optional, will create if not provided)
        include_metadata: Whether to include processing metadata in output
        
    Returns:
        Dictionary containing:
        - document_id: The document UUID used
        - file_name: Original file name
        - total_chunks: Number of chunks created
        - chunks: List of chunk dictionaries (if needed for further processing)
        
    Raises:
        FileNotFoundError: If the input file doesn't exist
        ValueError: If conversion or chunking fails
        Exception: If Weaviate operations fail
    """
    # Create client if not provided
    client_created = False
    if client is None:
        client = _get_weaviate_client()
        client_created = True
    
    try:
        # Step 1: Convert document to markdown using fast conversion
        logger.info(f"Fast processing document: {file_path}")

        markdown_content = fast_convert_to_markdown(
            file_path=file_path,
            include_metadata=include_metadata
        )
        
        if not markdown_content:
            raise ValueError(f"Failed to convert document to markdown: {file_path}")

        # Step 2: Chunk the markdown content
        chunks = chunk_markdown(
            markdown_input=markdown_content
        )

        if not chunks:
            raise ValueError(f"No chunks created from document: {file_path}")
        
        # Step 3: Save document metadata to Weaviate
        saved_doc_id = save_document_to_weaviate(
            client=client,
            file_path=file_path,
            chunks=chunks,
            document_id=document_id,
            tenant_id=tenant_id
        )
        
        # Step 4: Save chunks to Weaviate
        save_chunks_to_weaviate(
            client=client,
            chunks=chunks,
            document_id=saved_doc_id,
            tenant_id=tenant_id
        )
        
        # Return summary information
        result = {
            "document_id": saved_doc_id,
            "file_name": Path(file_path).name,
            "total_chunks": len(chunks),
            "chunks": chunks  # Include chunks in case caller needs them
        }
        
        logger.info(f"Successfully fast processed document '{Path(file_path).name}' with {len(chunks)} chunks")
        return result
        
    except FileNotFoundError as e:
        logger.error(f"File not found: {e}")
        raise
    except ValueError as e:
        logger.error(f"Processing error: {e}")
        raise
    except Exception as e:
        logger.error(f"Unexpected error processing document: {e}")
        raise
    finally:
        # Close client if we created it
        if client_created:
            try:
                client.close()
                logger.info("Closed Weaviate connection")
            except Exception as e:
                logger.warning(f"Error closing Weaviate connection: {e}")



